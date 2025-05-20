import os
import pickle
import time
import argparse
from google.auth.transport.requests import Request
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build, Resource
from googleapiclient.errors import HttpError
import io
from googleapiclient.http import MediaIoBaseDownload
import mimetypes

# External libraries for file processing
import docx  # for .docx
import PyPDF2  # for .pdf
from pptx import Presentation  # for .pptx
import pandas as pd  # for .xlsx
import openai
import json
import re
from ollama import chat

import os


# At the top of your file
OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
if not OPENAI_API_KEY:
    raise ValueError("OPENAI_API_KEY environment variable not set")

# Constants
SCOPES = [
    'https://www.googleapis.com/auth/drive.readonly',
    'https://www.googleapis.com/auth/documents.readonly',
    'https://www.googleapis.com/auth/presentations.readonly',
    'https://www.googleapis.com/auth/spreadsheets.readonly'
]

TOKEN_PICKLE = 'token.pickle'
CREDENTIALS_FILE = 'credentials.json'

# Supported file types
SUPPORTED_MIME_TYPES = {
    'application/vnd.google-apps.document',
    'application/vnd.google-apps.presentation',
    'application/vnd.google-apps.spreadsheet',
    'application/pdf',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    'text/plain'
}

SUPPORTED_EXTENSIONS = {
    '.txt', '.pdf', '.docx', '.pptx', '.xlsx', '.xls', '.csv'
}

# Framework and model configurations
FRAMEWORKS = {
    'openai': {
        'models': ['gpt-4o', 'gpt-4', 'gpt-3.5-turbo'],
        'function': lambda prompt, model: openai.OpenAI(api_key=OPENAI_API_KEY).chat.completions.create(
            model=model,
            messages=[{'role':'user','content':prompt}],
            temperature=0.3,
            max_tokens=600
        )
    },
    'ollama': {
        'models': ['qwen3:14b', 'llama3.1', 'llama3.2', 'deepseek-r1'],
        'function': lambda prompt, model: chat(
            model=model,
            messages=[{'role':'user','content':prompt}],
            options={
                "temperature": 0.3,
                'top_p': 0.95,
                'format': 'json',
                "stream": False,
                'num_predict': 4096
            }
        )
    }
}

def parse_arguments():
    """Parse command line arguments"""
    parser = argparse.ArgumentParser(description='Process Google Drive files and generate summaries.')
    parser.add_argument('--root-folder', required=True, help='Name of the root folder to process')
    parser.add_argument('--framework', choices=FRAMEWORKS.keys(), required=True, help='AI framework to use')
    parser.add_argument('--model', required=True, help='Model to use for summarization')
    parser.add_argument('--output', default='drive_summary.csv', help='Output CSV filename')
    return parser.parse_args()

def authenticate():
    """Authenticate the user and return valid credentials."""
    creds = None
    if os.path.exists(TOKEN_PICKLE):
        with open(TOKEN_PICKLE, 'rb') as token:
            creds = pickle.load(token)
    
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file(CREDENTIALS_FILE, SCOPES)
            creds = flow.run_local_server(port=0)
        
        with open(TOKEN_PICKLE, 'wb') as token:
            pickle.dump(creds, token)
    return creds

def connect_drive_api() -> Resource:
    """Connect to the Google Drive API and return the service resource."""
    creds = authenticate()
    return build('drive', 'v3', credentials=creds)

def is_supported_file(file_info: dict) -> bool:
    """Check if a file is of a supported type for processing."""
    mime_type = file_info.get('mimeType', '')
    name = file_info.get('name', '')
    
    # Check native Google file types
    if mime_type in SUPPORTED_MIME_TYPES:
        return True
    
    # Check file extensions for non-Google files
    ext = os.path.splitext(name)[1].lower()
    return ext in SUPPORTED_EXTENSIONS

def extract_text_from_google_file(file_id: str) -> str:
    """Extracts text from a Google Docs, Sheets, or Slides file using Google APIs."""
    creds = authenticate()
    
    try:
        drive_service = build('drive', 'v3', credentials=creds)
        metadata = drive_service.files().get(fileId=file_id, fields='mimeType, name').execute()
        mime_type = metadata['mimeType']
    except HttpError as e:
        print(f"Failed to retrieve metadata: {e}")
        return ""

    if mime_type == 'application/vnd.google-apps.document':
        docs_service = build('docs', 'v1', credentials=creds)
        doc = docs_service.documents().get(documentId=file_id).execute()
        return extract_text_from_doc_content(doc.get('body', {}).get('content', []))

    elif mime_type == 'application/vnd.google-apps.presentation':
        slides_service = build('slides', 'v1', credentials=creds)
        pres = slides_service.presentations().get(presentationId=file_id).execute()
        return extract_text_from_slides(pres.get('slides', []))

    elif mime_type == 'application/vnd.google-apps.spreadsheet':
        sheets_service = build('sheets', 'v4', credentials=creds)
        spreadsheet = sheets_service.spreadsheets().get(spreadsheetId=file_id).execute()
        sheet_titles = [s['properties']['title'] for s in spreadsheet.get('sheets', [])]

        result_text = []
        for title in sheet_titles:
            values = sheets_service.spreadsheets().values().get(
                spreadsheetId=file_id,
                range=title
            ).execute()
            rows = values.get('values', [])
            for row in rows:
                result_text.append('\t'.join(row))
        return '\n'.join(result_text)

    else:
        print(f"Unsupported Google file type: {mime_type}")
        return ""

def extract_text_from_doc_content(elements):
    """Extract text from Google Docs content elements"""
    lines = []
    for elem in elements:
        para = elem.get('paragraph')
        if para:
            for run in para.get('elements', []):
                text = run.get('textRun', {}).get('content', '')
                if text.strip():
                    lines.append(text.strip())
    return '\n'.join(lines)

def extract_text_from_slides(slides):
    """Extract text from Google Slides pages"""
    text_chunks = []
    for slide in slides:
        for element in slide.get('pageElements', []):
            shape = element.get('shape')
            if not shape:
                continue
            text_elements = shape.get('text', {}).get('textElements', [])
            for text_elem in text_elements:
                if 'textRun' in text_elem:
                    content = text_elem['textRun'].get('content', '')
                    if content.strip():
                        text_chunks.append(content.strip())
    return '\n'.join(text_chunks)

def download_file_content(file_id: str) -> str:
    """Download and parse a file from Google Drive into plain text."""
    drive_service = connect_drive_api()

    try:
        metadata = drive_service.files().get(
            fileId=file_id,
            fields="name, mimeType"
        ).execute()
        name = metadata['name']
        mime_type = metadata['mimeType']

        if not is_supported_file(metadata):
            print(f"[SKIP] Unsupported file type: {name} ({mime_type})")
            return ""

        fh = io.BytesIO()

        if mime_type.startswith("application/vnd.google-apps"):
            export_mime = {
                'application/vnd.google-apps.document': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                'application/vnd.google-apps.presentation': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                'application/vnd.google-apps.spreadsheet': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            }.get(mime_type)

            request = drive_service.files().export_media(fileId=file_id, mimeType=export_mime)
        else:
            request = drive_service.files().get_media(fileId=file_id)

        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            _, done = downloader.next_chunk()
        fh.seek(0)

    except HttpError as error:
        if error.resp.status == 403 and "exportSizeLimitExceeded" in str(error):
            # print(f"[SKIP] File too large to export: {name}")
            return extract_text_from_google_file(file_id)
        else:
            print(f"[ERROR] Failed to download file {name}: {error}")
        return ""

    # Parse the file content based on extension
    # ext = os.path.splitext(name)[1].lower()
    if mime_type == 'application/vnd.google-apps.document':
        ext = '.docx'
    elif mime_type == 'application/vnd.google-apps.presentation':
        ext = '.pptx'
    elif mime_type == 'application/vnd.google-apps.spreadsheet':
        ext = '.xlsx'
    else:
        ext = os.path.splitext(name)[1].lower()

    try:
        if ext == '.txt':
            return fh.read().decode('utf-8', errors='ignore')
        elif ext == '.pdf':
            reader = PyPDF2.PdfReader(fh)
            return '\n'.join(page.extract_text() or '' for page in reader.pages)
        elif ext == '.docx':
            doc = docx.Document(fh)
            return '\n'.join(p.text for p in doc.paragraphs)
        elif ext == '.pptx':
            prs = Presentation(fh)
            return '\n'.join(
                shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text')
            )
        elif ext in ['.xlsx', '.xls', '.csv']:
            df = pd.read_excel(fh, sheet_name=None) if ext != '.csv' else {'Sheet1': pd.read_csv(fh)}
            return '\n'.join(df[sheet].to_string() for sheet in df)
        else:
            print(f"[SKIP] Unsupported file type for parsing: {ext}")
            return ""
    except Exception as e:
        print(f"[ERROR] Failed to parse file {name}: {e}")
        return ""

def summarize_with_ai(text: str, framework: str, model: str) -> tuple:
    """Summarize text using the specified AI framework and model."""
    
    # categories = [
    #     'Life Skills', 'Dance', 'Disability Advocacy', 'Speech & Communication',
    #     'Expressive Therapy Activities', 'Healthy Relationships', 'Journaling'
    # ]
    
    categories = [
        "Plannings",
        "Proposals",
        "Agreements",
        "Project Status",
        "Client Presentations",
        "Results and Next Steps",
        "Meeting Minutes"
    ]

    prompt = (
        "Please summarize this text in no more than 100 words this is a must, just a quick summary, no analysis, "
        "the response should not be more than 100 words by any means. Then classify it into one of the following categories: "
        + ", ".join(categories) + ". "
        "Respond with a JSON object with keys 'summary' and 'category' use these exact words as i'm parsing on them.\n\n" + text
    )

    try:
        # Get the appropriate function based on framework
        ai_function = FRAMEWORKS[framework]['function']
        response = ai_function(prompt, model)
        
        # Handle different response formats
        if framework == 'openai':
            content = response.choices[0].message.content.strip()
        else:  # ollama
            content = response.message.content

        # Extract JSON from response
        match = re.search(r"\{.*\}", content, re.DOTALL)
        json_str = match.group(0) if match else content

        data = json.loads(json_str)
        return data.get('summary', ''), data.get('category', '')
    
    except json.JSONDecodeError:
        try:
            return data.get('Summary', ''), data.get('Category', '')
        except:
            print(f"❌ Failed to decode JSON response from {framework} {model}")
            return content, ''
    except Exception as e:
        print(f"❌ Error during summarization with {framework} {model}: {e}")
        return '', ''

def list_all_files(drive_service, folder_id: str, path: str = '') -> list:
    """Recursively list all supported files under a folder."""
    all_items = []
    query = f"'{folder_id}' in parents and trashed = false"
    page_token = None
    
    while True:
        response = drive_service.files().list(
            q=query,
            spaces='drive',
            fields="nextPageToken, files(id, name, mimeType, size, modifiedTime)",
            pageToken=page_token
        ).execute()
        
        for f in response.get('files', []):
            item_path = os.path.join(path, f['name'])
            if f['mimeType'] == 'application/vnd.google-apps.folder':
                all_items.extend(list_all_files(drive_service, f['id'], item_path))
            elif is_supported_file(f):
                f['path'] = item_path
                all_items.append(f)
        
        page_token = response.get('nextPageToken', None)
        if not page_token:
            break

    return all_items

def main():
    args = parse_arguments()
    
    # Validate model selection
    if args.model not in FRAMEWORKS[args.framework]['models']:
        print(f"Error: Model '{args.model}' is not supported by framework '{args.framework}'")
        print(f"Available models for {args.framework}: {', '.join(FRAMEWORKS[args.framework]['models'])}")
        return

    drive_service = connect_drive_api()
    
    # Find root folder ID
    folders = drive_service.files().list(
        q=f"mimeType='application/vnd.google-apps.folder' and name='{args.root_folder}' and trashed = false",
        spaces='drive',
        fields="files(id, name)",
        pageSize=10
    ).execute().get('files', [])

    if not folders:
        print(f"Folder '{args.root_folder}' not found.")
        return

    root_id = folders[0]['id']
    print(f"Processing files under '{args.root_folder}' (ID: {root_id}) using {args.framework}/{args.model}...")

    all_files = list_all_files(drive_service, root_id, args.root_folder)
    # print (all_files)
    rows = []
    
    for f in all_files:
        text = download_file_content(f['id'])
        if not text.strip():
            continue
            
        summary, category = summarize_with_ai(text, args.framework, args.model)
        rows.append({
            'path': f['path'],
            'name': f['name'],
            'category': category,
            'size': f.get('size', 'N/A'),
            'modified': f.get('modifiedTime', 'N/A'),
            'mimetype': f['mimeType'],
            'summary': summary,
        })
        print(f"Processed: {f['path']}")

    df = pd.DataFrame(rows)
    df.to_csv(args.output, index=False)
    print(f"Saved summary to {args.output}")

if __name__ == "__main__":
    main()